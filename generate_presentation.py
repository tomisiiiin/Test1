"""Python script to generate Project_Presentation.pptx from slide content.
Requires: python-pptx (pip install python-pptx)
Run: python generate_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt

slides = [
    {
        'title': 'Project Name',
        'lines': ['Group Members: Nafiur, Labib, Yixuan, Anya, Yi, Minhanjul, Tomisin', 'Date: 2025-10-22'],
        'notes': 'Welcome and quick one-line project description. Introduce the team.'
    },
    {
        'title': 'Agenda',
        'lines': ['Introduction & Brief Summary', 'Milestone 1: Project Setup', 'Milestone 2: Product Listing & Marketplace Operations', 'Demonstration Overview', 'Challenges, Tools, Next Steps, Q&A'],
        'notes': 'Walk through the agenda so the audience knows the flow.'
    },
    {
        'title': 'Milestone 1: Project Setup',
        'lines': ['Objective: Establish a functional, collaborative development environment', 'Key achievements: React & Node.js initialized; MongoDB connected via Mongoose; GitHub repo structured; Branching & PR workflow; Issue template and contribution guide'],
        'notes': 'Briefly explain why each achievement mattered.'
    },
    {
        'title': 'Sprint 1: Task Assignments',
        'lines': ['Frontend Setup — Nafiur', 'Backend Setup — Labib', 'Database Setup — Yixuan', 'User Models — Anya', 'Routes & Controllers — Yi', 'JWT Authentication — Labib', 'State Management — Nafiur', 'Backend Docs — Minhanjul', 'Frontend Docs — Tomisin'],
        'notes': 'Highlight how the work was divided to parallelize development.'
    },
    {
        'title': 'Milestone 2: Product Listing & Marketplace Operations',
        'lines': ['Objective: Build core functionality for listing, managing, and viewing products', 'Features: Product schema & CRUD APIs; Product listing UI; Profile management; User dashboard; Search feature'],
        'notes': 'Point out one or two screens you’ll demo.'
    },
    {
        'title': 'Progress Summary',
        'lines': ['Core backend APIs and DB connections: Complete', 'Frontend product listing & profiles: Functional', 'Authentication (JWT): Implemented', 'Repo documentation & contribution workflow: Finalized'],
        'notes': 'Summarize readiness: what’s stable and what’s next.'
    },
    {
        'title': 'Delivered Artifacts / Evidence',
        'lines': ['Environment setup guide', 'Repository documentation & Wiki', 'GitHub issue and milestone templates', 'Product list page (UI)', 'Login & Signup pages', 'Postman API tests (CRUD for products)'],
        'notes': 'Tell audience where to find the repo and docs.'
    },
    {
        'title': 'Demonstration Overview (Demo Flow)',
        'lines': ['Login/Signup', 'Create a product listing', 'View and edit listings', 'Edit user profile', 'Demonstrate search & navigation'],
        'notes': 'Announce you will switch to demo mode and any test creds.'
    },
    {
        'title': 'Challenges & Solutions',
        'lines': ['Coordinating parallel frontend/backend work', 'Authentication edge cases (token expiry)', 'UI/UX consistency across components', 'Solutions: PR rules & daily syncs; JWT refresh strategy; Shared UI component library'],
        'notes': 'Pick 1 challenge and describe the resolution.'
    },
    {
        'title': 'Tools & Technologies',
        'lines': ['Frontend: React', 'Backend: Node.js & Express', 'Database: MongoDB + Mongoose', 'Auth: JWT', 'Collaboration: GitHub, Postman, VS Code'],
        'notes': 'Briefly state why this stack was chosen.'
    },
    {
        'title': 'Next Steps / Roadmap',
        'lines': ['Short-term: Improve search UX & filters; Add image upload; Expand API tests and CI', 'Medium-term: Pagination & performance; Role-based access; Deployment to staging'],
        'notes': 'Give tentative timeline for next sprint and owners.'
    },
    {
        'title': 'Key Takeaways',
        'lines': ['Effective GitHub workflows accelerate collaboration', 'Practical MERN experience: auth, DB, UI integration', 'Strong documentation and testing reduce onboarding friction', 'Teamwork improved through clear assignments & PRs'],
        'notes': 'Summarize team learnings.'
    },
    {
        'title': 'Conclusion & Q&A',
        'lines': ['Thank you!', 'Repository: https://github.com/tomisiiiin/Test1', 'Contact: olaiyatomisin@gmail.com'],
        'notes': 'Invite questions and indicate next steps after Q&A.'
    }
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

for s in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = s['title']
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(s['lines']):
        if i == 0:
            p = body.paragraphs[0]
            p.text = line
            p.font.size = Pt(18)
        else:
            p = body.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(16)
    # add speaker notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = s['notes']

prs.save('Project_Presentation.pptx')